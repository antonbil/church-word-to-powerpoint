# sermon_core.py
from docx import Document
from pptx import Presentation
import os
from .sermon_extract import SermonExtract
from .sermon_create import SermonCreate
from .sermon_utils import SermonUtils
from .settings import Settings


class Sermon(SermonExtract, SermonCreate, SermonUtils):
    """
    Represents a church service and handles the creation of a PowerPoint presentation
    from a Word document order of service.
    """

    def __init__(self):
        """
        Initializes the Sermon object
        """
        # Get the directory of the current file (sermon_core.py)
        current_dir = os.path.dirname(os.path.abspath(__file__))
        # Get the name of the current directory
        current_dir_name = os.path.basename(current_dir)
        # Initialize the settings
        self.settings = Settings(current_dir)

        # Construct the full path to the word document
        self.word_filename = os.path.join(current_dir_name, self.settings.get_setting("default_word_filename"))
        self.powerpoint_template_filename = os.path.join(current_dir, self.settings.get_setting("powerpoint_template_filename"))
        self.max_reading_lines = self.settings.get_setting("max_reading_lines")
        self.powerpoint_filename = os.path.splitext(self.word_filename)[0] + ".pptx"
        self.word_document = None
        self.powerpoint_presentation = None
        # current_paragraph_index is the pointer to the current paragraph in the Word-file
        self.current_paragraph_index = 0
        self.num_paragraphs = 0
        # Define the tags
        self.tags = self.settings.get_tags()
        self.current_tag = None

    def load_word_document(self):
        """
        Loads the Word document and handles potential errors.
        """
        try:
            self.word_document = Document(self.word_filename)
        except FileNotFoundError:
            print(f"Error: Word document '{self.word_filename}' not found.")
            self.word_document = None
        except Exception as e:
            print(f"An unexpected error occurred while loading the Word document: {e}")
            self.word_document = None

    def create_powerpoint_presentation(self):
        """
        Creates an empty PowerPoint presentation with the specified filename.
        """
        try:
            # Load the template
            # Get the directory of the current file (sermon_core.py)
            current_dir = os.path.dirname(os.path.abspath(__file__))
            template_filename = os.path.join(current_dir, self.powerpoint_template_filename)
            self.powerpoint_presentation = Presentation(template_filename)
            self.powerpoint_presentation.save(self.powerpoint_filename)
        except Exception as e:
            print(f"An unexpected error occurred while creating the PowerPoint presentation: {e}")
            self.powerpoint_presentation = None

    def process_sermon(self):
        """
        Main method to process the sermon data and create the PowerPoint.
        Iterates through the paragraphs in the Word document using a while loop
        and a pointer, identifying start tags and calling the appropriate extraction
        functions.
        """
        self.load_word_document()
        if self.word_document is None:
            return  # Stop processing if there was an error loading the Word document

        self.create_powerpoint_presentation()
        if self.powerpoint_presentation is None:
            return

        paragraphs = self.word_document.paragraphs
        self.num_paragraphs = len(paragraphs)
        self.current_paragraph_index = 0

        # print(self.num_paragraphs)

        while self.current_paragraph_index < self.num_paragraphs:
            paragraph = self.word_document.paragraphs[self.current_paragraph_index]
            is_start_tag = False
            for tag_type, tag_data in self.tags.items():
                if tag_data["begin"] in paragraph.text:
                    is_start_tag = True

                    self.current_tag = tag_type

                    # Process the current section
                    if self.current_tag == "hymn":
                        title, hymn_data = self.extract_hymn_section(paragraphs[self.current_paragraph_index:])
                        self.create_hymn_slides(title, hymn_data)
                    elif self.current_tag == "offering":
                        offering_data = self.extract_offering_section(paragraphs[self.current_paragraph_index:])
                        self.create_offering_slides(offering_data)
                    elif self.current_tag == "intro":
                        intro_data = self.extract_intro_section(paragraphs[self.current_paragraph_index:])
                        self.current_paragraph_index += 1
                        self.create_intro_slides(intro_data, "slide-layout-intro-1")
                        self.create_intro_slides(intro_data, "slide-layout-intro-2")
                        self.create_intro_slides(intro_data, "slide-layout-intro-2", True)
                        self.create_empty_slide()
                    elif self.current_tag == "reading":
                        title, reading_data = self.extract_reading_section(paragraphs[self.current_paragraph_index:])
                        # print(reading_data)
                        self.create_reading_slides(title, reading_data)
                    elif self.current_tag == "outro":
                        date, parson, performed_piece = self.extract_outro_section(paragraphs[self.current_paragraph_index:])
                        # print(date, parson)
                        self.create_outro_slides(date, parson, "slide-layout-outro-1", performed_piece = performed_piece)
                        self.create_outro_slides(date, parson, "slide-layout-outro-2")
                    elif self.current_tag == "illustration":
                        image = self.extract_illustration(paragraphs[self.current_paragraph_index:])
                        self.create_illustration_slides(image)

            if not is_start_tag:
                self.current_paragraph_index += 1

        self.remove_slide(self.powerpoint_presentation)
        self.powerpoint_presentation.save(self.powerpoint_filename)
        if self.powerpoint_presentation is None:
            return

        print(f"PowerPoint presentation '{self.powerpoint_filename}' created successfully.")

    def remove_slide(self, prs):
        """
        Removes the first slide from a PowerPoint presentation.

        Args:
            prs: The Presentation object (pptx.presentation.Presentation) representing the PowerPoint presentation.
        """
        # Get a reference to the first slide
        selected_slide = prs.slides[0]

        # Get the slide_id of the slide to be removed
        slide_id = selected_slide.slide_id

        # Get the relationship ID (rId) of the first slide in the slide list
        # rIds is the Id of the xml in the presentation, that contains the info about the first slide
        rIds = prs.slides._sldIdLst[0].rId

        # Drop the relationship (remove the slide from the relationships)
        # This line removes the slide from the presentation's relationships
        prs.part.drop_rel(rIds)

        # Find the index of the slide to remove in the _sldIdLst
        # _sldIdLst is the internal list of slides in the presentation
        for idx, sld in enumerate(prs.slides._sldIdLst):
            if sld.id == slide_id:
                remove_idx = idx  # get the index of the slide

        # Remove the slide from the _sldIdLst
        # This line removes the slide from the internal slide list
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[remove_idx])