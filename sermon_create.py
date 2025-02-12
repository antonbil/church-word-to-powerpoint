# sermon_create.py
from docx.shared import Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt


class SermonCreate:
    """
    Contains the methods for creating PowerPoint slides.
    """
    def create_hymn_slides(self, title, hymn_data):
        """
        Creates PowerPoint slides for the hymn sections using a template.

        Args:
            title (str): The title of the hymn section (or None if no title).
            hymn_data (list): A list of dictionaries containing hymn data (text and images).
                          Each dictionary should have keys like "text" (str) and "images" (list).
        """
        print("add hymn-data")
        #create hymn-data to make code more clear
        image = None
        image_list = [h["images"][0]for h in hymn_data if len(h["images"]) > 0]
        if len(image_list) > 0:
            image = image_list[0]
        hymn_parts = [hymn["text"] for hymn in hymn_data if hymn["text"]]

        # Constants and settings
        template_id = "slide-layout-lied"
        song_length_first = self.settings.get_setting("powerpoint-hymn-song-length-first")
        song_length_first_image = self.settings.get_setting("powerpoint-hymn-song-length-first-image")
        song_length_rest = self.settings.get_setting("powerpoint-hymn-song-length-rest")

        # Early exit if no powerpoint.
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return
        # Early exit if there are no hymn parts.
        if not image and not hymn_parts:
            print("Warning: No hymn text and image found in hymn_data.")
            return

        #local function to get the number of lines in a string
        def get_number_lines(string):
            return len(string.split("\n"))

        original_template_id = template_id

        is_first_slide = True

        current_song_length = song_length_first
        current_length = 0

        if image:
            template_id = template_id + "-image"
            current_song_length = song_length_first_image

        for hymn_part in hymn_parts:
            # check if new slide is needed
            if get_number_lines(hymn_part) + current_length > current_song_length or is_first_slide:
                slide = self.add_slide(template_id)
                if not is_first_slide:
                    current_song_length = song_length_rest

            # Add title (only on the first slide)
            if title and is_first_slide:
                self.set_title(slide, title)

            is_first_slide = False
            template_id = original_template_id + "-no-title"

            # add image to slide if it exists
            if image:
                # Image creation and formatting
                # self._add_image_to_slide(image, slide)
                self.replace_image_in_placeholder(slide, image)
                image = None

            # add hymn-part to slide
            for i, placeholder in enumerate(slide.placeholders):
                if placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    if get_number_lines(hymn_part) + current_length > current_song_length:
                        placeholder.text = hymn_part
                    else:
                        placeholder.text = (placeholder.text + "\n\n" + hymn_part).strip()
                    current_length = get_number_lines(placeholder.text)

                    # Set content text appearance
                    for paragraph in placeholder.text_frame.paragraphs:
                        self.set_text_appearance(paragraph)
                    # set the text at the top:
                    placeholder.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        self.create_empty_slide()

    def create_reading_slides(self, title, reading_data):
        """
        Creates PowerPoint slides for the reading sections using a template.

        Args:
            title (str): The title of the reading section (or None if no title).
            reading_data (list): A list of dictionaries containing reading data (text and images).
        """
        print("add reading-data")
        template_id = "slide-layout-reading"
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return
        is_first_slide = True

        for reading in reading_data:
            slide = self.add_slide(template_id)

            # Add title (only on the first slide)
            if title and is_first_slide:
                self.set_title(slide, title)

                is_first_slide = False
                template_id = template_id + "-no-title"

            # find second placeholder

            for i, p in enumerate(slide.placeholders):
                if p.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    p.text = reading["text"]
                    # Set content text color to white
                    for paragraph in p.text_frame.paragraphs:
                        self.set_text_appearance(paragraph)
                    # set the text at the top:
                    p.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        self.create_empty_slide()

    def create_outro_slides(self, date, parson, template_id, performed_piece = None):
        """
        Creates the outro slide.

        Args:
            date (str): The date of the next service.
            parson (str): The name of the parson.
        """
        print("create_outro_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide = self.add_slide(template_id)

        # Set the title
        title_text = self.settings.get_setting("powerpoint-outro_title")
        if performed_piece:
            title_text = performed_piece

        def extra_layout_func(paragraph, line_number):
            paragraph.alignment = PP_ALIGN.CENTER

        self.set_title(slide, title_text, extra_layout_func)

        def content_layout_func(paragraph, line_number):
            paragraph.alignment = PP_ALIGN.CENTER
            font_size = self.settings.get_setting('powerpoint-outro_font_size')
            paragraph.font.size = Pt(font_size)
            paragraph.font.italic = True
            paragraph.font.bold = False

        # Set the content
        outro_data = {"parson":parson, "date":date}
        outro_template = self.settings.get_setting("powerpoint-outro_template")
        fields = ["date", "parson"]
        content_text = "\n".join(self.fill_template_with_data(fields, outro_data, outro_template))

        self.format_placeholder_text(1, content_text, slide, content_layout_func)

    def create_offering_slides(self, offering_data):
        """
        Creates the offering slide.

        Args:
            offering_data (list): The data of the offering (list).
        """
        print("create_offering_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide = self.add_slide("slide-layout-offering")

        # Set the content
        offering_template = self.settings.get_setting("powerpoint-offering_template")
        fields = ["offering_goal", "bank_account_number"]
        output = self.fill_template_with_data(fields, offering_data, offering_template)

        content_text = "\n".join(output)
        content_text_list = content_text.split("\n")
        empty_item = self.find_first_empty_string_index(content_text_list) + 2

        # add content to slide
        def extra_layout_func(paragraph, line_number):
            paragraph.font.underline = (True if line_number == 0 or line_number == empty_item else False)
            font_size = self.settings.get_setting('powerpoint-offering_font_size')
            paragraph.font.size = Pt(font_size)
            paragraph.font.italic = True
            paragraph.font.bold = False

        self.format_placeholder_text(1, content_text, slide, extra_layout_func)
        self.create_empty_slide()


    def create_intro_slides(self, intro_data, template_id, performed_piece_in_title = False):
        """
        Creates the intro slide.

        Args:
            intro_data (dict): The data of the intro (dict).
        """
        print("create_intro_slide")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide = self.add_slide(template_id)

        # Set the title
        title_text = self.settings.get_setting("powerpoint-intro_title")
        def extra_layout_func(paragraph, line_number):
            paragraph.alignment = PP_ALIGN.CENTER
        self.set_title(slide, title_text, extra_layout_func)

        performed_piece = ""
        intro_template = self.settings.get_setting('powerpoint-intro_template')

        fields = ["date", "parson", "theme", "organist"]
        intro_template = self.fill_template_with_data(fields, intro_data, intro_template)

        # Set the content
        intro_lines = intro_template
        if "performed_piece" in intro_data:
            performed_piece = intro_data['performed_piece']
        content_text = "\n".join(intro_lines)

        def extra_layout_func(paragraph, line_number):
            font_size = self.settings.get_setting('powerpoint-intro_font_size')
            paragraph.font.size = Pt(font_size)
            paragraph.font.italic = True
            paragraph.font.bold = False
        #add content to slide
        self.format_placeholder_text(1, content_text, slide, extra_layout_func)

        # set title to performed_piece
        if performed_piece and performed_piece_in_title:
            self.set_title(slide, performed_piece, extra_layout_func)

    def fill_template_with_data(self, fields, data_dict, template_list):
        """
        Fills a template list with data from a dictionary.

        This function iterates through a list of field names and attempts to replace
        corresponding placeholders within a template list with values from a data
        dictionary. If a field is not found in the data dictionary, the placeholder
        and potentially the preceding "heading-line" are removed from the template.

        Args:
            fields (list): A list of field names (strings) to look for in the template.
                           Example: ["date", "parson", "theme"]
            data_dict (dict): A dictionary where keys are field names and values are
                               the data to insert into the template.
                               Example: {"date": "2024-01-28", "parson": "John Doe"}
            template_list (list): A list of strings representing the template.
                                   Placeholders are denoted by curly braces (e.g., "{date}").
                                   Example: ["{date}", "", "Voorganger:", "{parson}"]

        Returns:
            list: The modified template list with placeholders replaced or removed.
        """
        for id in fields:
            try:
                # Try to find the index of the placeholder in the template
                index = template_list.index("{" + id + "}")

                # Check if the field is in intro_data and the index is valid
                if id in data_dict and index >= 0:
                    # Replace the placeholder with the data from intro_data
                    template_list[index] = template_list[index].replace("{" + id + "}", data_dict[id])
                else:
                    # Field not in intro_data or index invalid, so remove the placeholder
                    template_list.pop(index)  # Remove the placeholder

                    # Check if there is a previous line and it's not empty
                    if index > 0 and template_list[index - 1].strip() != "":
                        # Remove the heading-line if the previous line is not empty
                        template_list.pop(index - 1)
                    elif index > 0:
                        # remove the empty heading line
                        template_list.pop(index - 1)

            except ValueError:
                # The placeholder was not found in the template, so skip it
                pass

        return template_list

    def create_illustration_slides(self, image_data):
        """
        Creates a slide for the illustration section, displaying the image.

        Args:
            image_data: the image for the illustration-slide
        """
        print("create_illustration_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        if image_data is not None:
            slide = self.add_slide("slide-layout-intro-2")
            # Remove the title-placeholder
            for shp in slide.placeholders:
                if shp.name.startswith(self.settings.get_setting("placeholder-title-name")):
                    sp = shp.element
                    sp.getparent().remove(sp)
            self._add_image_to_slide(image_data, slide, setting_id = "illustration")

            self.create_empty_slide()

    def format_placeholder_text(self, placeholder_index, content_text, slide, custom_formatter=None):
        """
        Formats the text within a placeholder with the specified settings.

        Args:
            placeholder_index: the index of the placeholder.
            content_text: the text to put in the placeholder.
            slide: the current slide.
            custom_formatter: An optional callable (e.g., lambda) that takes a paragraph
                              and the line number as arguments for custom formatting.
        """
        for i, p in enumerate(slide.placeholders):
            if i == placeholder_index:
                p.text = content_text

                for line_number, paragraph in enumerate(p.text_frame.paragraphs):
                    self.set_text_appearance(paragraph)
                    # align the entire paragraph in the middle
                    paragraph.alignment = PP_ALIGN.CENTER
                    if custom_formatter:
                        custom_formatter(paragraph, line_number)

                # set the text at the top:
                p.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    def find_first_empty_string_index(self, string_list):
        """
        Finds the index of the first empty string in a list of strings.

        Args:
          string_list: A list of strings.

        Returns:
          The index of the first empty string, or None if no empty string is found.
        """
        for index, item in enumerate(string_list):
            if not item:
                return index
        return None  # Return None if no empty string is found

    def add_slide(self, slide_layout_code):
        """
        Add slide to the current powerpoint-presentation.

        Args:
            slide_layout (layout-from-template): The layout that is used.
        """

        slide_layout = self.settings.get_setting(slide_layout_code)

        # layout = None
        # for layout_idx, slide_layout1 in enumerate(self.powerpoint_presentation.slide_layouts):
        #     if layout_idx == slide_layout:
        #         print(f"\nSlide Layout {layout_idx}: {slide_layout1.name}")
        #         layout = slide_layout1

        slide = self.powerpoint_presentation.slides.add_slide(self.powerpoint_presentation.slide_layouts[slide_layout])
        return slide

    def create_empty_slide(self):
        """
        create empty slide based on default empty template-slide
        """
        self.add_slide("slide-layout-empty")

