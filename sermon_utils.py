# sermon_utils.py
import datetime
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import io
import re

class SermonUtils:
    """
    Contains utility methods for the Sermon class.
    """

    def format_date(self, date_text):
        """
        Formats the date from a string like "12-jan" to a string like "Zondag 12 januari 2025".

        Args:
            date_text (str): The date string in the format "dd-mmm" (e.g., "12-jan").

        Returns:
            str or None: The formatted date string (e.g., "Zondag 12 januari 2025") or None if the input date is invalid.
        """
        try:
            # Parse the input date string into a datetime object
            date_obj = datetime.datetime.strptime(date_text, "%d-%b")  # this function can raise a valueError exception

            # Replace the year with 2025 (this is a fixed year)
            date_obj = date_obj.replace(year=2025)

            # Get the Dutch day of the week name (e.g., "Zondag")
            day_of_week = self.get_day_of_week(date_obj.weekday())

            # Format the date into the desired string format
            formatted_date = date_obj.strftime("%d %B %Y")

            # Combine the day of the week and the formatted date
            return f"{day_of_week} {formatted_date}"

        except ValueError:
            # Handle the case where the date string is in an invalid format
            return None

    def get_day_of_week(self, weekday_number):
        """
        Returns the Dutch day of the week name for a given weekday number (0-6).

        Args:
            weekday_number (int): The weekday number (0 for Monday, 1 for Tuesday, ..., 6 for Sunday).

        Returns:
            str or None: The Dutch day of the week name (e.g., "Maandag", "Dinsdag", etc.) or None if the weekday number is invalid.
        """
        # List of Dutch day of the week names
        days = ["Maandag", "Dinsdag", "Woensdag", "Donderdag", "Vrijdag", "Zaterdag", "Zondag"]

        # Check if the weekday number is within the valid range (0-6)
        if 0 <= weekday_number <= 6:
            return days[weekday_number]  # Return the day name from the list
        else:
            return None  # Return None if the weekday number is invalid
    def check_end_tag(self, tag, paragraph):
        """
        Checks if any of the defined end tags for the given tag are present in the paragraph.

        Args:
            tag (str): The tag to check (e.g., "hymn", "reading").
            paragraph (docx.paragraph.Paragraph): The paragraph to check.

        Returns:
            bool: True if any of the end tags are found, False otherwise.
        """
        for end_tag in self.tags[tag]["end"]:
            if end_tag.lower() in paragraph.text.lower():
                return True
        return False

    def get_hymn_title(self, index, paragraphs):
        """
        Extracts the title of a hymn from a list of paragraphs.

        Args:
            self: The instance of the class.
            index (int): The current index in the list of paragraphs.
                         This may be updated if the title spans multiple paragraphs.
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects)
                               representing the potential hymn section.

        Returns:
            tuple: A tuple containing:
                - in_hymn_section (bool): True if the section is identified as a hymn section, False otherwise.
                - index (int): The updated index after processing the title paragraphs.
                - title (str): The extracted title of the hymn (stripped of whitespace), or None if no title is found.
        """
        title = None  # Initialize the title to None (no title found yet)
        in_hymn_section = False  # Initially, we are not in a hymn section

        first_paragraph = paragraphs[0]  # Get the first paragraph of the potential hymn section

        # Check if the first paragraph contains a title (indicated by bold text)
        if first_paragraph.runs and first_paragraph.runs[0].bold:
            title = first_paragraph.text  # Extract the text of the first paragraph as the potential title

            # Check if the title contains the hymn section start tag
            if self.tags["hymn"]["begin"] in title:
                title = title.replace(self.tags["hymn"]["begin"], "")  # Remove the start tag from the title
                in_hymn_section = True  # Mark that we are now in the hymn section

            # Check if this is a two-line title (title spans over two paragraphs)
            if len(paragraphs) > 1:
                second_paragraph = paragraphs[1]  # Get the second paragraph

                # Check if the second paragraph is also part of the title (indicated by bold text)
                if second_paragraph.runs and second_paragraph.runs[0].bold:
                    title = title + "\n" + second_paragraph.text  # Append the text of the second paragraph to the title, separated by a newline
                    index += 1  # Increment the index because we've processed an extra paragraph

        # Return the results
        # in_hymn_section: indicates if we've identified a hymn section
        # index: the updated index (increased if a two-line title was found)
        # title.strip(): the extracted title (with leading/trailing whitespace removed), or None
        return in_hymn_section, index, title.strip()

    def get_reading_title(self, index, paragraphs):
        """
        Extracts the title of a reading from a list of paragraphs, and skips leading empty paragraphs.

        Args:
            self: The instance of the class.
            index (int): The current index in the list of paragraphs.
                         This may be updated if the title spans multiple paragraphs or if empty paragraphs are skipped.
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects)
                               representing the potential reading section.

        Returns:
            tuple: A tuple containing:
                - in_reading_section (bool): True if the section is identified as a reading section, False otherwise.
                - index (int): The updated index after processing the title and skipping empty paragraphs.
                - title (str): The extracted title of the reading (stripped of whitespace), or None if no title is found.
        """
        title = None  # Initialize the title to None (no title found yet)
        in_reading_section = False  # Initially, we are not in a reading section
        first_paragraph = paragraphs[0]  # Get the first paragraph
        start_check_empty_paragraphs = 1  # Initialize the variable that indicates which empty paragraph must be checked.

        # Check if the first paragraph contains a title (indicated by bold text)
        if first_paragraph.runs and first_paragraph.runs[0].bold:
            title = first_paragraph.text  # Extract the text of the first paragraph as the potential title

            # Check if the title contains the reading section start tag
            if self.tags["reading"]["begin"] in title:
                title = title.replace(self.tags["reading"]["begin"], "")  # Remove the start tag from the title
                in_reading_section = True  # Mark that we are now in the reading section
                index += 1  # increment the index

            title = title.strip()  # Remove leading/trailing whitespace from the title

            # Check if this is a two-line title (title spans over two paragraphs)
            if len(paragraphs) > 1:
                second_paragraph = paragraphs[1]  # Get the second paragraph

                # Check if the second paragraph is also part of the title (indicated by bold text)
                if second_paragraph.runs and second_paragraph.runs[0].bold:
                    title = title + "\n" + second_paragraph.text.strip()  # Append the text of the second paragraph to the title, separated by a newline (also strip this paragraph)
                    index += 1  # Increment the index because we've processed an extra paragraph
                    start_check_empty_paragraphs = 2  # set the index to the right value

        # Skip empty paragraphs at the top of the reading section (if there is no title or two-line title)
        while (start_check_empty_paragraphs < len(paragraphs) and paragraphs[start_check_empty_paragraphs].text == ""
               and not self.paragraph_content_contains_image(paragraphs[start_check_empty_paragraphs])):
            start_check_empty_paragraphs += 1  # Move to the next paragraph
            index += 1  # Increment the index because we've skipped a paragraph

        # Return the results
        # in_reading_section: indicates if we've identified a reading section
        # index: the updated index (increased if a two-line title was found or if empty paragraphs were skipped)
        # title.strip(): the extracted title (with leading/trailing whitespace removed), or None
        return in_reading_section, index, title.strip()

    def _extract_image_embeds(self, paragraph):
        """
        Helper function to extract image embeds from a paragraph.

        Args:
            paragraph: The paragraph object (docx.paragraph.Paragraph).

        Returns:
            list: A list of image embed IDs (strings).
        """
        embeds = []
        for run in paragraph.runs:
            for drawing in run._element.xpath('.//w:drawing'):
                for inline in drawing.xpath('.//wp:inline'):
                    for graphic in inline.xpath('.//a:graphic'):
                        for graphicData in graphic.xpath('.//a:graphicData'):
                            for pic in graphicData.xpath('.//pic:pic'):
                                for blipfill in pic.xpath('.//pic:blipFill'):
                                    for blip in blipfill.xpath('.//a:blip'):
                                        embed = blip.get(
                                            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                        if embed:
                                            embeds.append(embed)
        return embeds

    def extract_paragraph_content(self, paragraph):
        """
        Extracts the text content and images (as bytes) from a paragraph.

        Args:
            paragraph: The paragraph object (docx.paragraph.Paragraph).

        Returns:
            dict: A dictionary containing the paragraph text and a list of image bytes.
        """
        paragraph_data = {"text": paragraph.text, "images": []}
        embeds = self._extract_image_embeds(paragraph)
        for embed in embeds:
            image_part = self.word_document.part.related_parts[embed]
            image_bytes = image_part.blob
            paragraph_data["images"].append(image_bytes)
        return paragraph_data

    def paragraph_content_contains_image(self, paragraph):
        """
        Checks if a paragraph contains any images.

        Args:
            paragraph: The paragraph object (docx.paragraph.Paragraph).

        Returns:
            bool: True if the paragraph contains at least one image, False otherwise.
        """
        embeds = self._extract_image_embeds(paragraph)
        return len(embeds) > 0

    def calculate_text_height(self, text, font_size):
        """
        Calculates an approximate height for a given text block and font size.
        """
        lines = text.count('\n') + 1  # Count lines based on newline characters
        line_height = font_size * 1.2  # Approximate line height
        total_height_inches = (lines * line_height) / 72  # Convert points to inches
        return Inches(total_height_inches)

    def duplicate_slide_with_layout(self, slide, layout_number, extra_layout_func=None):
        """
        Duplicates a slide and applies a specific layout to the new slide,
        copying title and content (text and pictures) to the corresponding placeholders.
        Uses the provided extra_layout_func to apply formatting.

        Args:
            slide: The slide to duplicate.
            layout_number (int): The index of the slide layout to apply (0-based).
            extra_layout_func: A function to apply extra layout settings to the title.

        Returns:
            The duplicated slide with the specified layout, or None if an error occurred.
        """
        prs = slide.part.package.presentation_part.presentation

        # Check if the layout index is valid
        if layout_number < 0 or layout_number >= len(prs.slide_layouts):
            print(f"Error: Invalid layout index {layout_number}. Must be between 0 and {len(prs.slide_layouts) - 1}.")
            return None

        # Get the desired slide layout
        desired_layout = prs.slide_layouts[layout_number]

        # Create a new slide with the desired layout
        new_slide = prs.slides.add_slide(desired_layout)
        # for shp1 in new_slide.placeholders:
        #     print("new ph", shp1.name)
        content_placeholder_id = self.settings.get_setting("placeholder-content-name")

        # Dictionary to map placeholder types
        placeholder_map = {
            "title": 0,  # Placeholder id for title in the default layout
            "body": 1   # Placeholder id for content in the default layout
        }

        # Copy content from the original slide to the new slide
        # print("search", content_placeholder_id)
        for shp in slide.placeholders:
            if shp.name.startswith("Title"):
                # Found title placeholder in original slide
                try:
                    # Copy title text to the new slide using extra_layout_func
                    title_text = shp.text
                    self.set_title(new_slide, title_text, extra_layout_func)

                except Exception as e:
                    print(f"An error occurred copying the title: {e}")
                    pass

            if shp.name.startswith(content_placeholder_id):
                # print(shp.name)
                # print("tekst:")
                # print(shp.text)
                # Found content placeholder in original slide

                try:
                    # Get the body text frame in the new slide
                    # print("1")
                    self.format_placeholder_text(1, shp.text, new_slide)

                    # Copy pictures from original slide to new slide
                    for inner_shp in shp.shapes:
                        try:
                            if inner_shp.has_picture:
                                picture = inner_shp.picture
                                left, top, width, height = picture.left, picture.top, picture.width, picture.height
                                new_slide.shapes.add_picture(picture.image.blob, left, top, width, height)
                        except Exception as e:
                            #print(f"An error occurred copying the picture: {e}")
                            pass

                except Exception as e:
                    print(f"An error occurred copying the content: {e}")
                    pass
            try:
                if shp.has_table:
                    tbl = shp.table
                    print(f"table: {tbl}")
            except Exception as e:
                print(f"An error occurred: {e}")
                pass
            try:
                if shp.has_chart:
                    chart = shp.chart
                    print(f"chart: {chart}")
            except Exception as e:
                print(f"An error occurred: {e}")
                pass
            try:
                if shp.has_picture:
                    picture = shp.picture
                    print(f"picture: {picture}")
                    # Add picture to new slide
                    left, top, width, height = picture.left, picture.top, picture.width, picture.height
                    new_slide.shapes.add_picture(picture.image.blob, left, top, width, height)
            except Exception as e:
                # print(f"An error occurred: {e}")
                pass

        return new_slide

    def _add_image_to_slide(self, image_data, slide, setting_id = "hymn"):
        """Adds an image to a slide with border and shadow.

        This function adds an image to the specified slide using the given image data
        and the image settings (width, height, left, top) defined in the settings.
        It also adds a border and a shadow to the image.

        Args:
            image_data (bytes): The image data in bytes.
            slide (pptx.slide.Slide): The slide object to which the image will be added.
        Returns:
            image: The image created, or None if error.
        """
        # Get image dimensions and position from settings
        image_width = Inches(self.settings.get_setting(setting_id + "-image_width"))  # Get the image width in inches
        image_height = Inches(self.settings.get_setting(setting_id + "-image_height"))  # Get the image height in inches
        image_left = Inches(self.settings.get_setting(setting_id + "-image_left"))  # Get the image left position in inches
        image_top = Inches(self.settings.get_setting(setting_id + "-image_top"))  # Get the image top position in inches

        # Convert image data to BytesIO object for PowerPoint
        image_stream = io.BytesIO(image_data)

        try:
            # Add the image to the slide using the specified dimensions and position
            picture = slide.shapes.add_picture(image_stream, left=image_left, top=image_top, width=image_width, height=image_height)

            # Add a rectangle shape behind the picture to create a border and shadow effect
            self.add_border_and_shadow(slide, picture)
            return picture

        except Exception as e:
            # Handle any exceptions that occur during image addition
            print(f"An error occurred while adding the picture: {e}")
            return None

    def add_border_and_shadow(self, slide, picture):
        """Adds a border and a shadow to the picture.

        This function adds a rectangle shape behind the picture to create a border and shadow effect.

        Args:
            slide: The slide object to which the border and shadow will be added.
            picture: The picture object to which the border and shadow will be added.
        """
        # Get the position and size of the image
        left, top, width, height = picture.left, picture.top, picture.width, picture.height

        # Calculate the position and size of the border shape
        border_width = width + Inches(0.02)  # 0.02 inches = 2 * 0.01, to add 1 pixel to both left and right
        border_height = height + Inches(0.02)  # 0.02 inches = 2 * 0.01, to add 1 pixel to both top and bottom
        border_left = left - Inches(0.01)  # 0.01 inches, to shift 1 pixel to the left
        border_top = top - Inches(0.01)  # 0.01 inches, to shift 1 pixel up

        # Add a rectangle shape behind the picture
        border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, border_left, border_top, border_width, border_height)

        # set the order of the shapes, so that the image is placed on top of the shape
        # The shape is placed BEFORE the picture in the XML, so the shape is behind the picture
        picture.element.addprevious(border_shape.element)

        # Set the border (line) properties
        line = border_shape.line
        # Set the color of the line in the border
        # get color "powerpoint-image-border_color" from settings
        color = self.settings.get_setting("powerpoint-image-border_color")
        # RGBColor expects red, green, and blue components as integers between 0 and 255
        line.color.rgb = RGBColor(color["red"], color["green"], color["blue"])
        line.width = 12700  # 1 pt (1pt = 12700 emu)
        # Add a shadow to the shape
        # Set the shadow
        shadow = border_shape.shadow
        #shadow.inherit = True
        #print(shadow)
        # Create a shadow

        # shadow.visible

    def replace_image_in_placeholder(self, slide, image_data):
        """Replaces the content of an image placeholder on a slide with a new image.

        Args:
            slide (pptx.slide.Slide): The slide containing the image placeholder.
            image_data (bytes): The new image data in bytes.
        """
        # find image-placeholder
        image_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                image_placeholder = placeholder
                break

        if not image_placeholder:
            print("Error: No image placeholder found on the slide.")
            return

        # Convert image data to BytesIO object for PowerPoint
        image_stream = io.BytesIO(image_data)

        # Add image to placeholder.
        image_placeholder.insert_picture(image_stream)

    def set_title(self, slide, title_text, custom_formatter=None):
        """
        Sets the title text of a slide and applies custom formatting.

        Args:
            slide: The slide object (pptx.slide.Slide) to set the title on.
            title_text (str): The text to set as the slide title.
            custom_formatter (function, optional): An optional function that takes two arguments:
                - paragraph (pptx.text.text._Paragraph): The paragraph object of the title.
                - line_number (int): The line number of the paragraph.
                This function can be used to apply custom formatting to the title paragraph.
                Defaults to None.
        """
        try:
            title_placeholder = slide.shapes.title  # Get the title placeholder shape on the slide

            # Set the title text
            title_placeholder.text = title_text  # Set the text of the title placeholder to the given title_text

            # Loop through each paragraph in the title's text frame
            for line_number, paragraph in enumerate(title_placeholder.text_frame.paragraphs):
                self.set_color_text_line(paragraph) #Set the color of the text.
                # Set the font size of the paragraph
                paragraph.font.size = Pt(self.settings.get_setting("powerpoint-title_font_size")) # Set the font size based on the settings
                # set the font type
                paragraph.font.name = self.settings.get_setting("powerpoint-title_font_type")  # set the font-type
                paragraph.font.italic = True
                paragraph.font.bold = True

                # Apply custom formatting if a custom formatter function is provided
                if custom_formatter:
                    custom_formatter(paragraph, line_number)  # Call the custom formatter function, passing the paragraph and line number

        except Exception as e:
            print(e)  # Print any exceptions that occur during title setting

    def set_text_appearance(self, paragraph):
        """
        Sets the text appearance (color, size, and font type) of a paragraph.

        Args:
            paragraph: The paragraph object (pptx.text.text._Paragraph) to format.
        """
        # Set the text color of the paragraph
        self.set_color_text_line(paragraph)

        # Set the font size of the paragraph
        paragraph.font.size = Pt(self.settings.get_setting("powerpoint-content_font_size"))

        # Set the font type (font name) of the paragraph
        paragraph.font.name = self.settings.get_setting("powerpoint-content_font_type")  # set the font-type
        paragraph.font.bold = True
        paragraph.level = 0

        # Access the paragraph formatting properties using _pPr (paragraph properties)
        pPr = paragraph._pPr
        # Remove bullet point if it exists
        if pPr.find(qn('a:buNone')) is None and pPr.find(qn('a:buAutoNum')) is None:
            # Create the buNone element and add it to the pPr
            buNone = pPr.makeelement(qn('a:buNone'))
            pPr.append(buNone)

    def split_string_list(self, string_list):
        """Splits a list of strings into sublists based on empty strings or numbered lines.

        Args:
            string_list (list): A list of strings.

        Returns:
            list: A list of lists of strings.
        """
        result = []
        current_sublist = []

        for item in string_list:
            # Check if the string is empty or starts with a number followed by a space or period
            if not item or re.match(r"^\d+[\s.]", item):
                # If the current sublist is not empty, add it to the result
                if current_sublist:
                    result.append(current_sublist)
                # Start a new sublist
                current_sublist = []
            # Add the current item to the current sublist
            if item:
                current_sublist.append(item)

        # Add the last sublist if it's not empty
        if current_sublist:
            result.append(current_sublist)

        return result

    def set_color_text_line(self, paragraph):
        """
        Sets the color of the text in a paragraph.

        Args:
            paragraph: The paragraph object (pptx.text.text._Paragraph) whose text color should be changed.
        """
        # Get the color settings from the configuration
        color = self.settings.get_setting("powerpoint-content_font_color")

        # Set the color of the text in the paragraph
        # RGBColor expects red, green, and blue components as integers between 0 and 255
        paragraph.font.color.rgb = RGBColor(color["red"], color["green"], color["blue"])