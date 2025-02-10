from docx import Document
from docx.shared import Pt, Inches
from pptx import Presentation
from pptx.util import Inches
import os
import io
import re
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from .settings import Settings
import datetime

class Sermon:
    """
    Represents a church service and handles the creation of a PowerPoint presentation
    from a Word document order of service.
    """

    def __init__(self):
        """
        Initializes the Sermon object
        """
        # Get the directory of the current file (sermon_core.py)
        current_dir = os.path.dirname(os.path.abspath(__file__))
        # Get the name of the current directory
        current_dir_name = os.path.basename(current_dir)
        # Initialize the settings
        self.settings = Settings(current_dir)

        # Construct the full path to the word document
        self.word_filename = os.path.join(current_dir_name, self.settings.get_setting("default_word_filename"))
        self.powerpoint_template_filename = os.path.join(current_dir, self.settings.get_setting("powerpoint_template_filename"))
        self.max_reading_lines = self.settings.get_setting("max_reading_lines")
        self.powerpoint_filename = os.path.splitext(self.word_filename)[0] + ".pptx"
        self.word_document = None
        self.powerpoint_presentation = None
        # current_paragraph_index is the pointer to the current paragraph in the Word-file
        self.current_paragraph_index = 0
        self.num_paragraphs = 0
        # Define the tags
        self.tags = self.settings.get_tags()
        self.current_tag = None

    def load_word_document(self):
        """
        Loads the Word document and handles potential errors.
        """
        try:
            self.word_document = Document(self.word_filename)
        except FileNotFoundError:
            print(f"Error: Word document '{self.word_filename}' not found.")
            self.word_document = None
        except Exception as e:
            print(f"An unexpected error occurred while loading the Word document: {e}")
            self.word_document = None

    def create_powerpoint_presentation(self):
        """
        Creates an empty PowerPoint presentation with the specified filename.
        """
        try:
            # Load the template
            # Get the directory of the current file (sermon_core.py)
            current_dir = os.path.dirname(os.path.abspath(__file__))
            template_filename = os.path.join(current_dir, self.powerpoint_template_filename)
            self.powerpoint_presentation = Presentation(template_filename)
            self.powerpoint_presentation.save(self.powerpoint_filename)
        except Exception as e:
            print(f"An unexpected error occurred while creating the PowerPoint presentation: {e}")
            self.powerpoint_presentation = None

    def process_sermon(self):
        """
        Main method to process the sermon data and create the PowerPoint.
        Iterates through the paragraphs in the Word document using a while loop
        and a pointer, identifying start tags and calling the appropriate extraction
        functions.
        """
        self.load_word_document()
        if self.word_document is None:
            return  # Stop processing if there was an error loading the Word document

        self.create_powerpoint_presentation()
        if self.powerpoint_presentation is None:
            return

        paragraphs = self.word_document.paragraphs
        self.num_paragraphs = len(paragraphs)
        self.current_paragraph_index = 0

        # print(self.num_paragraphs)

        while self.current_paragraph_index < self.num_paragraphs:
            paragraph = self.word_document.paragraphs[self.current_paragraph_index]
            is_start_tag = False
            for tag_type, tag_data in self.tags.items():
                if tag_data["begin"] in paragraph.text:
                    is_start_tag = True

                    self.current_tag = tag_type

                    # Process the current section
                    if self.current_tag == "hymn":
                        title, hymn_data = self.extract_hymn_section(paragraphs[self.current_paragraph_index:])
                        self.create_hymn_slides(title, hymn_data)
                    elif self.current_tag == "offering":
                        offering_data = self.extract_offering_section(paragraphs[self.current_paragraph_index:])
                        self.create_offering_slides(offering_data)
                    elif self.current_tag == "intro":
                        intro_data = self.extract_intro_section(paragraphs[self.current_paragraph_index:])
                        self.current_paragraph_index += 1
                        self.create_intro_slides(intro_data)
                    elif self.current_tag == "reading":
                        title, reading_data = self.extract_reading_section(paragraphs[self.current_paragraph_index:])
                        # print(reading_data)
                        self.create_hymn_slides(title, reading_data)
                    elif self.current_tag == "outro":
                        date, parson, performed_piece = self.extract_outro_section(paragraphs[self.current_paragraph_index:])
                        # print(date, parson)
                        self.create_outro_slides(date, parson, performed_piece)
                        self.create_outro_slides(date, parson)
                    elif self.current_tag == "illustration":
                        image = self.extract_illustration(paragraphs[self.current_paragraph_index:])
                        self.create_illustration_slides(image)

            if not is_start_tag:
                self.current_paragraph_index += 1

        self.remove_slide(self.powerpoint_presentation)
        self.powerpoint_presentation.save(self.powerpoint_filename)
        if self.powerpoint_presentation is None:
            return

        print(f"PowerPoint presentation '{self.powerpoint_filename}' created successfully.")

    def remove_slide(self, prs):
        """
        Removes the first slide from a PowerPoint presentation.

        Args:
            prs: The Presentation object (pptx.presentation.Presentation) representing the PowerPoint presentation.
        """
        # Get a reference to the first slide
        selected_slide = prs.slides[0]

        # Get the slide_id of the slide to be removed
        slide_id = selected_slide.slide_id

        # Get the relationship ID (rId) of the first slide in the slide list
        # rIds is the Id of the xml in the presentation, that contains the info about the first slide
        rIds = prs.slides._sldIdLst[0].rId

        # Drop the relationship (remove the slide from the relationships)
        # This line removes the slide from the presentation's relationships
        prs.part.drop_rel(rIds)

        # Find the index of the slide to remove in the _sldIdLst
        # _sldIdLst is the internal list of slides in the presentation
        for idx, sld in enumerate(prs.slides._sldIdLst):
            if sld.id == slide_id:
                remove_idx = idx  # get the index of the slide

        # Remove the slide from the _sldIdLst
        # This line removes the slide from the internal slide list
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[remove_idx])

    def format_date(self, date_text):
        """
        Formats the date from a string like "12-jan" to a string like "Zondag 12 januari 2025".

        Args:
            date_text (str): The date string in the format "dd-mmm" (e.g., "12-jan").

        Returns:
            str or None: The formatted date string (e.g., "Zondag 12 januari 2025") or None if the input date is invalid.
        """
        try:
            # Parse the input date string into a datetime object
            date_obj = datetime.datetime.strptime(date_text, "%d-%b")  # this function can raise a valueError exception

            # Replace the year with 2025 (this is a fixed year)
            date_obj = date_obj.replace(year=2025)

            # Get the Dutch day of the week name (e.g., "Zondag")
            day_of_week = self.get_day_of_week(date_obj.weekday())

            # Format the date into the desired string format
            formatted_date = date_obj.strftime("%d %B %Y")

            # Combine the day of the week and the formatted date
            return f"{day_of_week} {formatted_date}"

        except ValueError:
            # Handle the case where the date string is in an invalid format
            return None

    def get_day_of_week(self, weekday_number):
        """
        Returns the Dutch day of the week name for a given weekday number (0-6).

        Args:
            weekday_number (int): The weekday number (0 for Monday, 1 for Tuesday, ..., 6 for Sunday).

        Returns:
            str or None: The Dutch day of the week name (e.g., "Maandag", "Dinsdag", etc.) or None if the weekday number is invalid.
        """
        # List of Dutch day of the week names
        days = ["Maandag", "Dinsdag", "Woensdag", "Donderdag", "Vrijdag", "Zaterdag", "Zondag"]

        # Check if the weekday number is within the valid range (0-6)
        if 0 <= weekday_number <= 6:
            return days[weekday_number]  # Return the day name from the list
        else:
            return None  # Return None if the weekday number is invalid
    def check_end_tag(self, tag, paragraph):
        """
        Checks if any of the defined end tags for the given tag are present in the paragraph.

        Args:
            tag (str): The tag to check (e.g., "hymn", "reading").
            paragraph (docx.paragraph.Paragraph): The paragraph to check.

        Returns:
            bool: True if any of the end tags are found, False otherwise.
        """
        for end_tag in self.tags[tag]["end"]:
            if end_tag.lower() in paragraph.text.lower():
                return True
        return False

    def get_hymn_title(self, index, paragraphs):
        """
        Extracts the title of a hymn from a list of paragraphs.

        Args:
            self: The instance of the class.
            index (int): The current index in the list of paragraphs.
                         This may be updated if the title spans multiple paragraphs.
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects)
                               representing the potential hymn section.

        Returns:
            tuple: A tuple containing:
                - in_hymn_section (bool): True if the section is identified as a hymn section, False otherwise.
                - index (int): The updated index after processing the title paragraphs.
                - title (str): The extracted title of the hymn (stripped of whitespace), or None if no title is found.
        """
        title = None  # Initialize the title to None (no title found yet)
        in_hymn_section = False  # Initially, we are not in a hymn section

        first_paragraph = paragraphs[0]  # Get the first paragraph of the potential hymn section

        # Check if the first paragraph contains a title (indicated by bold text)
        if first_paragraph.runs and first_paragraph.runs[0].bold:
            title = first_paragraph.text  # Extract the text of the first paragraph as the potential title

            # Check if the title contains the hymn section start tag
            if self.tags["hymn"]["begin"] in title:
                title = title.replace(self.tags["hymn"]["begin"], "")  # Remove the start tag from the title
                in_hymn_section = True  # Mark that we are now in the hymn section

            # Check if this is a two-line title (title spans over two paragraphs)
            if len(paragraphs) > 1:
                second_paragraph = paragraphs[1]  # Get the second paragraph

                # Check if the second paragraph is also part of the title (indicated by bold text)
                if second_paragraph.runs and second_paragraph.runs[0].bold:
                    title = title + "\n" + second_paragraph.text  # Append the text of the second paragraph to the title, separated by a newline
                    index += 1  # Increment the index because we've processed an extra paragraph

        # Return the results
        # in_hymn_section: indicates if we've identified a hymn section
        # index: the updated index (increased if a two-line title was found)
        # title.strip(): the extracted title (with leading/trailing whitespace removed), or None
        return in_hymn_section, index, title.strip()

    def get_reading_title(self, index, paragraphs):
        """
        Extracts the title of a reading from a list of paragraphs, and skips leading empty paragraphs.

        Args:
            self: The instance of the class.
            index (int): The current index in the list of paragraphs.
                         This may be updated if the title spans multiple paragraphs or if empty paragraphs are skipped.
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects)
                               representing the potential reading section.

        Returns:
            tuple: A tuple containing:
                - in_reading_section (bool): True if the section is identified as a reading section, False otherwise.
                - index (int): The updated index after processing the title and skipping empty paragraphs.
                - title (str): The extracted title of the reading (stripped of whitespace), or None if no title is found.
        """
        title = None  # Initialize the title to None (no title found yet)
        in_reading_section = False  # Initially, we are not in a reading section
        first_paragraph = paragraphs[0]  # Get the first paragraph
        start_check_empty_paragraphs = 1  # Initialize the variable that indicates which empty paragraph must be checked.

        # Check if the first paragraph contains a title (indicated by bold text)
        if first_paragraph.runs and first_paragraph.runs[0].bold:
            title = first_paragraph.text  # Extract the text of the first paragraph as the potential title

            # Check if the title contains the reading section start tag
            if self.tags["reading"]["begin"] in title:
                title = title.replace(self.tags["reading"]["begin"], "")  # Remove the start tag from the title
                in_reading_section = True  # Mark that we are now in the reading section
                index += 1  # increment the index

            title = title.strip()  # Remove leading/trailing whitespace from the title

            # Check if this is a two-line title (title spans over two paragraphs)
            if len(paragraphs) > 1:
                second_paragraph = paragraphs[1]  # Get the second paragraph

                # Check if the second paragraph is also part of the title (indicated by bold text)
                if second_paragraph.runs and second_paragraph.runs[0].bold:
                    title = title + "\n" + second_paragraph.text.strip()  # Append the text of the second paragraph to the title, separated by a newline (also strip this paragraph)
                    index += 1  # Increment the index because we've processed an extra paragraph
                    start_check_empty_paragraphs = 2  # set the index to the right value

        # Skip empty paragraphs at the top of the reading section (if there is no title or two-line title)
        while (start_check_empty_paragraphs < len(paragraphs) and paragraphs[start_check_empty_paragraphs].text == ""
               and not self.paragraph_content_contains_image(paragraphs[start_check_empty_paragraphs])):
            start_check_empty_paragraphs += 1  # Move to the next paragraph
            index += 1  # Increment the index because we've skipped a paragraph

        # Return the results
        # in_reading_section: indicates if we've identified a reading section
        # index: the updated index (increased if a two-line title was found or if empty paragraphs were skipped)
        # title.strip(): the extracted title (with leading/trailing whitespace removed), or None
        return in_reading_section, index, title.strip()

    def _extract_image_embeds(self, paragraph):
        """
        Helper function to extract image embeds from a paragraph.

        Args:
            paragraph: The paragraph object (docx.paragraph.Paragraph).

        Returns:
            list: A list of image embed IDs (strings).
        """
        embeds = []
        for run in paragraph.runs:
            for drawing in run._element.xpath('.//w:drawing'):
                for inline in drawing.xpath('.//wp:inline'):
                    for graphic in inline.xpath('.//a:graphic'):
                        for graphicData in graphic.xpath('.//a:graphicData'):
                            for pic in graphicData.xpath('.//pic:pic'):
                                for blipfill in pic.xpath('.//pic:blipFill'):
                                    for blip in blipfill.xpath('.//a:blip'):
                                        embed = blip.get(
                                            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                        if embed:
                                            embeds.append(embed)
        return embeds

    def extract_paragraph_content(self, paragraph):
        """
        Extracts the text content and images (as bytes) from a paragraph.

        Args:
            paragraph: The paragraph object (docx.paragraph.Paragraph).

        Returns:
            dict: A dictionary containing the paragraph text and a list of image bytes.
        """
        paragraph_data = {"text": paragraph.text, "images": []}
        embeds = self._extract_image_embeds(paragraph)
        for embed in embeds:
            image_part = self.word_document.part.related_parts[embed]
            image_bytes = image_part.blob
            paragraph_data["images"].append(image_bytes)
        return paragraph_data

    def paragraph_content_contains_image(self, paragraph):
        """
        Checks if a paragraph contains any images.

        Args:
            paragraph: The paragraph object (docx.paragraph.Paragraph).

        Returns:
            bool: True if the paragraph contains at least one image, False otherwise.
        """
        embeds = self._extract_image_embeds(paragraph)
        return len(embeds) > 0

    def calculate_text_height(self, text, font_size):
        """
        Calculates an approximate height for a given text block and font size.
        """
        lines = text.count('\n') + 1  # Count lines based on newline characters
        line_height = font_size * 1.2  # Approximate line height
        total_height_inches = (lines * line_height) / 72  # Convert points to inches
        return Inches(total_height_inches)

    def duplicate_slide_with_layout(self, slide, layout_number, extra_layout_func=None):
        """
        Duplicates a slide and applies a specific layout to the new slide,
        copying title and content (text and pictures) to the corresponding placeholders.
        Uses the provided extra_layout_func to apply formatting.

        Args:
            slide: The slide to duplicate.
            layout_number (int): The index of the slide layout to apply (0-based).
            extra_layout_func: A function to apply extra layout settings to the title.

        Returns:
            The duplicated slide with the specified layout, or None if an error occurred.
        """
        prs = slide.part.package.presentation_part.presentation

        # Check if the layout index is valid
        if layout_number < 0 or layout_number >= len(prs.slide_layouts):
            print(f"Error: Invalid layout index {layout_number}. Must be between 0 and {len(prs.slide_layouts) - 1}.")
            return None

        # Get the desired slide layout
        desired_layout = prs.slide_layouts[layout_number]

        # Create a new slide with the desired layout
        new_slide = prs.slides.add_slide(desired_layout)
        # for shp1 in new_slide.placeholders:
        #     print("new ph", shp1.name)
        content_placeholder_id = self.settings.get_setting("placeholder-content-name")

        # Dictionary to map placeholder types
        placeholder_map = {
            "title": 0,  # Placeholder id for title in the default layout
            "body": 1   # Placeholder id for content in the default layout
        }

        # Copy content from the original slide to the new slide
        # print("search", content_placeholder_id)
        for shp in slide.placeholders:
            if shp.name.startswith("Title"):
                # Found title placeholder in original slide
                try:
                    # Copy title text to the new slide using extra_layout_func
                    title_text = shp.text
                    self.set_title(new_slide, title_text, extra_layout_func)

                except Exception as e:
                    print(f"An error occurred copying the title: {e}")
                    pass

            if shp.name.startswith(content_placeholder_id):
                # print(shp.name)
                # print("tekst:")
                # print(shp.text)
                # Found content placeholder in original slide

                try:
                    # Get the body text frame in the new slide
                    # print("1")
                    self.format_placeholder_text(1, shp.text, new_slide)

                    # Copy pictures from original slide to new slide
                    for inner_shp in shp.shapes:
                        try:
                            if inner_shp.has_picture:
                                picture = inner_shp.picture
                                left, top, width, height = picture.left, picture.top, picture.width, picture.height
                                new_slide.shapes.add_picture(picture.image.blob, left, top, width, height)
                        except Exception as e:
                            #print(f"An error occurred copying the picture: {e}")
                            pass

                except Exception as e:
                    print(f"An error occurred copying the content: {e}")
                    pass
            try:
                if shp.has_table:
                    tbl = shp.table
                    print(f"table: {tbl}")
            except Exception as e:
                print(f"An error occurred: {e}")
                pass
            try:
                if shp.has_chart:
                    chart = shp.chart
                    print(f"chart: {chart}")
            except Exception as e:
                print(f"An error occurred: {e}")
                pass
            try:
                if shp.has_picture:
                    picture = shp.picture
                    print(f"picture: {picture}")
                    # Add picture to new slide
                    left, top, width, height = picture.left, picture.top, picture.width, picture.height
                    new_slide.shapes.add_picture(picture.image.blob, left, top, width, height)
            except Exception as e:
                # print(f"An error occurred: {e}")
                pass

        return new_slide

    def _add_image_to_slide(self, image_data, slide):
        """Adds an image to a slide.

        This function adds an image to the specified slide using the given image data
        and the image settings (width, height, left, top) defined in the settings.

        Args:
            image_data (bytes): The image data in bytes.
            slide (pptx.slide.Slide): The slide object to which the image will be added.
        """
        # Get image dimensions and position from settings
        image_width = Inches(self.settings.get_setting("powerpoint-image_width"))  # Get the image width in inches
        image_height = Inches(self.settings.get_setting("powerpoint-image_height"))  # Get the image height in inches
        image_left = Inches(self.settings.get_setting("powerpoint-image_left"))  # Get the image left position in inches
        image_top = Inches(self.settings.get_setting("powerpoint-image_top"))  # Get the image top position in inches

        # Convert image data to BytesIO object for PowerPoint
        image_stream = io.BytesIO(image_data)

        try:
            # Add the image to the slide using the specified dimensions and position
            slide.shapes.add_picture(image_stream, left=image_left, top=image_top, width=image_width, height=image_height)
        except Exception as e:
            # Handle any exceptions that occur during image addition
            print(f"An error occurred while adding the picture: {e}")

    def set_title(self, slide, title_text, custom_formatter=None):
        """
        Sets the title text of a slide and applies custom formatting.

        Args:
            slide: The slide object (pptx.slide.Slide) to set the title on.
            title_text (str): The text to set as the slide title.
            custom_formatter (function, optional): An optional function that takes two arguments:
                - paragraph (pptx.text.text._Paragraph): The paragraph object of the title.
                - line_number (int): The line number of the paragraph.
                This function can be used to apply custom formatting to the title paragraph.
                Defaults to None.
        """
        try:
            title_placeholder = slide.shapes.title  # Get the title placeholder shape on the slide

            # Set the title text
            title_placeholder.text = title_text  # Set the text of the title placeholder to the given title_text

            # Loop through each paragraph in the title's text frame
            for line_number, paragraph in enumerate(title_placeholder.text_frame.paragraphs):
                self.set_color_text_line(paragraph) #Set the color of the text.
                # Set the font size of the paragraph
                paragraph.font.size = Pt(self.settings.get_setting("powerpoint-title_font_size")) # Set the font size based on the settings

                # Apply custom formatting if a custom formatter function is provided
                if custom_formatter:
                    custom_formatter(paragraph, line_number)  # Call the custom formatter function, passing the paragraph and line number

        except Exception as e:
            print(e)  # Print any exceptions that occur during title setting

    def split_string_list(self, string_list):
        """Splits a list of strings into sublists based on empty strings or numbered lines.

        Args:
            string_list (list): A list of strings.

        Returns:
            list: A list of lists of strings.
        """
        result = []
        current_sublist = []

        for item in string_list:
            # Check if the string is empty or starts with a number followed by a space or period
            if not item or re.match(r"^\d+[\s.]", item):
                # If the current sublist is not empty, add it to the result
                if current_sublist:
                    result.append(current_sublist)
                # Start a new sublist
                current_sublist = []
            # Add the current item to the current sublist
            if item:
                current_sublist.append(item)

        # Add the last sublist if it's not empty
        if current_sublist:
            result.append(current_sublist)

        return result

    def set_color_text_line(self, paragraph):
        """
        Sets the color of the text in a paragraph.

        Args:
            paragraph: The paragraph object (pptx.text.text._Paragraph) whose text color should be changed.
        """
        # Get the color settings from the configuration
        color = self.settings.get_setting("powerpoint-content_font_color")

        # Set the color of the text in the paragraph
        # RGBColor expects red, green, and blue components as integers between 0 and 255
        paragraph.font.color.rgb = RGBColor(color["red"], color["green"], color["blue"])

    def extract_hymn_section(self, paragraphs):
        """
        Extracts the hymn section, including title, and all hymns within the section
        (text and images) from a list of paragraphs.

        Args:
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects).

        Returns:
            tuple: (title, hymn_data)
                   title (str or None): The title of the hymn section (or None if no title).
                   hymn_data (list): A list of dictionaries, where each dictionary represents a
                                     hymn and contains its text and image data.
        """
        hymn_data = []
        title = None
        index = -1
        print("check for hymn sections")

        # check if the hymn-section has a title:
        if len(paragraphs) > 0:
            in_hymn_section, index, title = self.get_hymn_title(index, paragraphs)
        outro_data = {
            "image": None
        }
        def add_image_function(image, outro_data):
            outro_data["image"] = image

        def add_line_function(paragraph, current_text, _):
            current_text.append(paragraph.text)

        text = self._extract_section_text("hymn", add_image_function=add_image_function,
                                   add_line_function=add_line_function, outro_data=outro_data).strip()
        if outro_data["image"]:
            paragraph_data = {"text": "", "images": [outro_data["image"]]}
            hymn_data.append(paragraph_data)

        text = self.remove_title_from_text(text, title)

        split_list = self.split_string_list(text.split("\n"))
        for hymn in split_list:
            if len(hymn) == 0:
                continue
            paragraph_data = {"text": "\n".join(hymn), "images": []}
            hymn_data.append(paragraph_data)

        return title, hymn_data

    def remove_title_from_text(self, text, title):
        """
            Removes the first line of text if it starts with the title.

            Args:
                text (str): The text to potentially modify.
                title (str): The title to check against.

            Returns:
                str: The modified text, or the original text if it didn't start with the title.
            """
        # workaround because the title of the reading sometimes is repeated as the first line of the content
        if text.startswith(title):
            # Find the first newline character
            first_newline_index = text.find('\n')

            if first_newline_index != -1:
                # Remove the first line (including the newline)
                text = text[first_newline_index + 1:]
            else:
                text = text[len(title):]
        return text

    def extract_offering_section(self, paragraphs):
        """
        Extracts the offering section (text) from a list of paragraphs.

        Args:
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects).

        Returns:
            tuple: (offering_data)
                   offering_data (list): A list of dictionaries, where each dictionary
                                        represents a part of the offering and contains
                                        its text.
        """
        print("extract_offering_section")
        def add_line_function(paragraph, current_text, _):
            cleaned_line = re.sub(r' {5,}', '\n', paragraph.text.strip())

            current_text.append(cleaned_line)
        full_text = self. _extract_section_text("offering", add_line_function = add_line_function)
        offering_goal, bank_account_number = self.extract_bank_account_number(full_text)
        offering_data = {"offering_goal": offering_goal, "bank_account_number": bank_account_number}
        return offering_data

    def _extract_section_text(self, section_name, add_line_function = None, outro_data = None, add_image_function = None):
        """
        Extracts text from a specified section in the Word document.

        Args:
            section_name (str): The name of the section to extract text from (e.g., "offering").

        Returns:
            list: A list of strings, where each string is a cleaned line of text from the section.
            int: the index of the last processed paragraph.
        """
        current_text = []
        in_section = False
        new_index = 0
        index = -1

        while self.current_paragraph_index + index < self.num_paragraphs - 1:
            index += 1
            paragraph = self.word_document.paragraphs[self.current_paragraph_index + index]

            if self.tags[section_name]["begin"] in paragraph.text:
                # A new section has started
                in_section = True
                t = paragraph.text.split(self.tags[section_name]["begin"])[-1]
                if t and add_line_function:
                    paragraph.text = t
                    add_line_function(paragraph, current_text, outro_data)
                continue

            if self.check_end_tag(section_name, paragraph):
                # End tag found for the section
                in_section = False
                new_index = index
                break

            if len(paragraph.text.strip()) == 0 and not in_section:
                # Empty line outside the section, skip it
                new_index = index
                continue

            if len(paragraph.text.strip()) > 0 and not in_section:
                # Text found outside the section, so the section is finished
                new_index = index
                break

            if in_section:
                if len(paragraph.text) > 0 and add_line_function:
                    add_line_function(paragraph, current_text, outro_data)

                new_index = index
                if add_image_function and self.paragraph_content_contains_image(paragraph):
                    add_image_function(self.extract_paragraph_content(paragraph)["images"][0], outro_data)
        current_text = "\n".join(current_text)
        self.current_paragraph_index = self.current_paragraph_index + new_index

        return current_text


    def extract_intro_section(self, paragraphs):
        """
        Extracts the introduction section (date, time, parson, theme, organist) from a list of paragraphs.

        Args:
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects).

        Returns:
            dict: A dictionary containing the extracted information:
                  {
                      "date": "Zondag 5 januari 2025",
                      "time": "10.00 uur",
                      "parson": "ds. Elly van Kuijk-Spaans",
                      "theme": "“Zaaien: een gids voor beginners”",
                      "organist": "Martin van der Bent"
                  }
        """
        print("extract_intro_section")
        intro_data = {}
        current_text = []

        def add_line_function(paragraph, current_text, _):
            current_text.append(paragraph.text.strip())
        intro_text = self. _extract_section_text("intro", add_line_function = add_line_function)

        sermon = self.settings.get_setting('word-intro-date_label')
        # Extract date
        sermon_date_list = [l for l in current_text if sermon in l]
        date_text = "25 december 2024"
        if len(sermon_date_list) > 0:
            l2 = sermon_date_list[0].split(sermon)
            if len(l2) > 1:
                date_text = l2[1].strip()

        # Convert month names to numbers
        month_numbers = {
            "januari": "1",
            "februari": "2",
            "maart": "3",
            "april": "4",
            "mei": "5",
            "juni": "6",
            "juli": "7",
            "augustus": "8",
            "september": "9",
            "oktober": "10",
            "november": "11",
            "december": "12",
        }

        intro_data["date"] = date_text
        if date_text:
            day, month_name, year = date_text.split()
            month_number = month_numbers.get(month_name.lower())
            if month_number:
                numeric_date_str = f"{day} {month_number} {year}"
                try:
                    date_object = datetime.strptime(numeric_date_str, "%d %m %Y")
                    intro_data["date"] = f"{self.get_day_of_week(date_object.weekday())} {date_text}"
                except ValueError as e:
                    print(f"Error parsing date: {e}")

        # Extract time
        time_match = re.search(r"(\d{1,2}\.\d{2}\suur)", intro_text)
        if time_match:
            intro_data["time"] = time_match.group(1)

        # Extract parson
        parson_text = self.settings.get_setting("word-intro-parson_text")
        parson_match = re.search(rf"{parson_text}\s*\n\s*(.+)", intro_text)
        if parson_match:
            intro_data["parson"] = parson_match.group(1).strip()

        # Regex to find "Thema:" and capture the text on the following line(s)
        theme_text = self.settings.get_setting("word-intro-theme_text")
        theme_match = re.search(rf"{theme_text}\s*“([^“”]+)”", intro_text)
        if theme_match:
            theme = theme_match.group(1).strip()
            intro_data["theme"] = theme

        # Extract organ-player and performed piece
        organ_text = self.settings.get_setting("word-intro-organplayer_text")
        organist_match = re.search(rf"{organ_text}\s+(.+):\s+(.+)", intro_text)
        if organist_match:
            intro_data["organist"] = organist_match.group(1).strip()
            intro_data["performed_piece"] = organist_match.group(2).strip()

        # print(intro_data)
        return intro_data

    def extract_reading_section(self, paragraphs):
        """
        Extracts the reading section from a list of paragraphs.

        Args:
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects).

        Returns:
            tuple: (title, reading_data)
                   title (str or None): The title of the reading section (or None if no title).
                   reading_data (list): A list of dictionaries, where each dictionary
                                        represents a part of the reading and contains
                                        its text (and potentially images, though they're not expected here).
        """
        print("extract_reading_section")
        reading_data = []
        title = None
        index = -1
        # check if the reading-section has a title:
        if len(paragraphs) > 0:
            in_reading_section, index, title = self.get_reading_title(index, paragraphs)

        def add_line_function(paragraph, current_text, _):
            cleaned_line = re.sub(r' {5,}', '\n', paragraph.text.strip())
            current_text.append(cleaned_line)
        full_text = self. _extract_section_text("reading", add_line_function = add_line_function)
        # workaround because the title of the reading sometimes is repeated as the first line of the content
        full_text = self.remove_title_from_text(full_text, title).strip()

        lines = full_text.split('\n')

        if len(lines) > self.max_reading_lines:
            parts = [lines[i:i + self.max_reading_lines] for i in range(0, len(lines), self.max_reading_lines)]
            for part in parts:
                reading_data.append({"text": "\n".join(part), "images": []})
        else:
            reading_data.append({"text": full_text, "images": []})
        return title, reading_data

    def extract_illustration(self, paragraphs):
        """
        Extracts the illustration from the given paragraphs.

        Args:
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects).

        Returns:
            tuple: (image_data, image_content_type) or (None, None) if no illustration is found.
        """
        print("extract_illustration")

        outro_data = {
            "image": None
        }
        def add_image_function(image, outro_data):
            outro_data["image"] = image

        self._extract_section_text("illustration", add_image_function=add_image_function, outro_data=outro_data)
        return outro_data["image"]


    def extract_outro_section(self, paragraphs):
        """
        Extracts the date and parson from the outro section.

        Args:
            paragraphs (list): A list of paragraphs (docx.paragraph.Paragraph objects).

        Returns:
            tuple: (date, parson) or (None, None) if not found.
        """
        print("extract_outro_section")
        current_text = []

        # Use a dictionary to store the values, that will be used as the closure for the function add_line_function
        outro_data = {
            "date_text": "",
            "parson": "",
            "performed_piece": "",
            "previous_line_is_sermon": False,
            "organ_text": self.settings.get_setting("word-outro-organ_text"),
            "next_sermon_text": self.settings.get_setting("word-outro-next_sermon_text")
        }

        # define the function, that will use the closure that is defined in the lines above
        def add_line_function(paragraph, _, outro_data):
            if outro_data["organ_text"] in paragraph.text:
                organ_text = outro_data["organ_text"]
                performed_piece_match = re.search(rf"{organ_text}\s*(.+)", paragraph.text.strip()) # use f-string
                if performed_piece_match:
                    outro_data["performed_piece"] = performed_piece_match.group(1).strip()
            if outro_data["next_sermon_text"].lower() in paragraph.text.lower():
                outro_data["previous_line_is_sermon"] = True
                return

            if outro_data["previous_line_is_sermon"]:
                parts = paragraph.text.split('\t')
                outro_data["previous_line_is_sermon"] = False
                if len(parts) >= 2:
                    date_text1 = parts[0]
                    parson1 = parts[len(parts) - 1]
                    date_text1 = self.format_date(date_text1)
                    outro_data["date_text"] = date_text1
                    outro_data["parson"] = parson1
        self._extract_section_text("outro", add_line_function = add_line_function, outro_data = outro_data)
        #make sure the program will end here
        self.current_paragraph_index = 200000

        return outro_data["date_text"], outro_data["parson"], outro_data["performed_piece"]

    def extract_bank_account_number(self, text):
        """
        Extracts the offering goal and bank account number from the text.

        Args:
            text (str): The text from the offering section.

        Returns:
            tuple: (offering_goal, bank_account_number)
        """
        bank_account_number = ""
        offering_goal = ""
        # Split the text at the "blauwe zak"
        blue_bag_text = self.settings.get_setting("word-offering-blue_bag_text")
        parts = text.split(blue_bag_text)
        first_part_of_offering_text = parts[0]

        # Regex to find the bank account number (IBAN) in the first part
        bank_account_regex = "([A-Z]{2}\d{2}\s[A-Z]{4}\s\d{4}\s\d{4}\s\d{2})"
        bank_number_match = re.search(bank_account_regex, first_part_of_offering_text)
        if bank_number_match:
            bank_account_number = bank_number_match.group(1).strip()

        # Clean the bank account number
        bank_account_number = re.sub(r"[^a-zA-Z0-9]", " ", bank_account_number)  # Replace non-alphanumeric with space
        bank_account_number = re.sub(r"\s{2,}", " ", bank_account_number)  # Replace multiple spaces with one
        bank_account_number = bank_account_number.strip()

        # Regex to find offering goal (text after "1ste (rode zak) Diaconie:")
        red_bag_text = self.settings.get_setting("word-offering-red_bag_text")
        diaconie_text = self.settings.get_setting("word-offering-diaconie_text")

        offering_goal_regex = fr"1ste \({red_bag_text}\)\s*{diaconie_text}\s*(.*?)(?=\n|$)"
        goal_match = re.search(offering_goal_regex, text)
        if goal_match:
            offering_goal = goal_match.group(1).strip()
        return offering_goal, bank_account_number


    def create_hymn_slides(self, title, hymn_data):
        """
        Creates PowerPoint slides for the hymn sections using a template.

        Args:
            title (str): The title of the hymn section (or None if no title).
            hymn_data (list): A list of dictionaries containing hymn data (text and images).
        """
        print("add hymn-data")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide_layout = self.powerpoint_presentation.slide_layouts[0]  # Assuming you want to use the first slide layout from the template
        is_first_slide = True

        last_bottom = 0
        previous_text = None
        for hymn in hymn_data:
            if not (last_bottom > 0 and hymn["text"]) and not (previous_text and hymn["text"]):
                slide = self.add_slide(slide_layout)

            # Add title (only on the first slide)
            if title and is_first_slide:
                self.set_title(slide, title)

                is_first_slide = False
            # add content (only image or text)
            if len(hymn["images"]) > 0:
                # Image formatting
                self._add_image_to_slide(hymn["images"][0], slide)

            elif hymn["text"]:
                i = 0
                for p in slide.placeholders:
                    if i==1:
                        if last_bottom > 0:
                            left = p.left
                            p.top = p.top + last_bottom - 200
                            p.left = left
                            previous_text = None
                            p.text = hymn["text"]
                        else:
                            if previous_text:
                                p.text = p.text + "\n\n" + hymn["text"]
                                previous_text = None
                            else:
                                p.text = hymn["text"]
                                previous_text = p
                        # Set content text color to white
                        for paragraph in p.text_frame.paragraphs:
                            self.set_color_text_line(paragraph)
                            paragraph.font.size = Pt(self.settings.get_setting("powerpoint-content_font_size"))
                        # set the text at the top:
                        p.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                        last_bottom = 0
                    i = i + 1
        self.create_empty_slide()

    def create_outro_slides(self, date, parson, performed_piece = None):
        """
        Creates the outro slide.

        Args:
            date (str): The date of the next service.
            parson (str): The name of the parson.
        """
        print("create_outro_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide_layout = self.powerpoint_presentation.slide_layouts[0]
        slide = self.add_slide(slide_layout)

        # Set the title
        title_text = self.settings.get_setting("powerpoint-outro_title")
        if performed_piece:
            title_text = performed_piece
        def extra_layout_func(paragraph, line_number):
            paragraph.alignment = PP_ALIGN.CENTER

        self.set_title(slide, title_text, extra_layout_func)

        # Set the content
        next_service_line = self.settings.get_setting("powerpoint-outro_next_service_line")
        parson_line = self.settings.get_setting("powerpoint-outro_parson_line") #new
        content_text = f"{next_service_line}:\n\n{date} \n\n{parson_line}:\n{parson}"
        self.format_placeholder_text(1, content_text, slide)

    def create_offering_slides(self, offering_data):
        """
        Creates the offering slide.

        Args:
            offering_data (list): The data of the offering (list).
        """
        print("create_offering_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide_layout = self.powerpoint_presentation.slide_layouts[0]
        slide = self.add_slide(slide_layout)

        # Set the content
        first_goal_label = self.settings.get_setting("powerpoint-offering_first_goal_label")
        second_goal_label = self.settings.get_setting("powerpoint-offering_second_goal_label")
        output = [first_goal_label, offering_data['offering_goal' ], offering_data["bank_account_number" ], "\n", second_goal_label]
        content_text = "\n".join(output)
        content_text_list = content_text.split("\n")
        empty_item = self.find_first_empty_string_index(content_text_list) + 2

        # add content to slide
        def extra_layout_func(paragraph, line_number):
            paragraph.font.underline = (True if line_number == 0 or line_number == empty_item else False)
        self.format_placeholder_text(1, content_text, slide, extra_layout_func)
        self.create_empty_slide()


    def create_intro_slides(self, intro_data):
        """
        Creates the intro slide.

        Args:
            intro_data (dict): The data of the intro (dict).
        """
        print("create_intro_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        slide_layout = self.powerpoint_presentation.slide_layouts[0]
        slide = self.add_slide(slide_layout)

        # Set the title
        title_text = self.settings.get_setting("powerpoint-intro_title")
        def extra_layout_func(paragraph, line_number):
            paragraph.alignment = PP_ALIGN.CENTER
        self.set_title(slide, title_text, extra_layout_func)

        performed_piece = ""
        # Set the content
        intro_lines = []
        if "date" in intro_data:
            intro_lines.append(f"{intro_data['date']}")
        if "parson" in intro_data:
            intro_lines.append(f"\n{self.settings.get_setting('powerpoint-outro_parson_line')}\n{intro_data['parson']}\n")
        if "theme" in intro_data:
            intro_lines.append(f"\n{self.settings.get_setting('powerpoint-intro_theme_line')}\n\"{intro_data['theme']}\"\n")
        if "organist" in intro_data:
            intro_lines.append(f"\n{self.settings.get_setting('powerpoint-intro_organist_line')}\n{intro_data['organist']}")
        if "performed_piece" in intro_data:
            performed_piece = intro_data['performed_piece']
        content_text = "\n".join(intro_lines)

        #add content to slide
        self.format_placeholder_text(1, content_text, slide)

        self.duplicate_slide_with_layout(slide, self.settings.get_setting('slide-layout-intro-2'), extra_layout_func)
        third_slide = self.duplicate_slide_with_layout(slide, self.settings.get_setting('slide-layout-intro-2'), extra_layout_func)
        # set title to performed_piece
        if performed_piece:
            self.set_title(third_slide, performed_piece, extra_layout_func)
        self.create_empty_slide()

    def create_illustration_slides(self, image_data):
        """
        Creates a slide for the illustration section, displaying the image.

        Args:
            image_data: the image for the illustration-slide
        """
        print("create_illustration_slides")
        if not self.powerpoint_presentation:
            print("Error: PowerPoint presentation not initialized.")
            return

        if image_data is not None:
            slide_layout = self.powerpoint_presentation.slide_layouts[0]
            slide = self.add_slide(slide_layout)
            # Remove the title-placeholder
            for shp in slide.placeholders:
                if shp.name.startswith(self.settings.get_setting("placeholder-title-name")):
                    sp = shp.element
                    sp.getparent().remove(sp)
            self._add_image_to_slide(image_data, slide)

            self.create_empty_slide()

    def format_placeholder_text(self, placeholder_index, content_text, slide, custom_formatter=None):
        """
        Formats the text within a placeholder with the specified settings.

        Args:
            placeholder_index: the index of the placeholder.
            content_text: the text to put in the placeholder.
            slide: the current slide.
            custom_formatter: An optional callable (e.g., lambda) that takes a paragraph
                              and the line number as arguments for custom formatting.
        """
        for i, p in enumerate(slide.placeholders):
            if i == placeholder_index:
                p.text = content_text

                for line_number, paragraph in enumerate(p.text_frame.paragraphs):
                    self.set_color_text_line(paragraph)
                    paragraph.font.size = Pt(self.settings.get_setting("powerpoint-content_font_size"))
                    # align the entire paragraph in the middle
                    paragraph.alignment = PP_ALIGN.CENTER
                    if custom_formatter:
                        custom_formatter(paragraph, line_number)

                # set the text at the top:
                p.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    def find_first_empty_string_index(self, string_list):
        """
        Finds the index of the first empty string in a list of strings.

        Args:
          string_list: A list of strings.

        Returns:
          The index of the first empty string, or None if no empty string is found.
        """
        for index, item in enumerate(string_list):
            if not item:
                return index
        return None  # Return None if no empty string is found

    def add_slide(self, slide_layout):
        """
        Add slide to the current powerpoint-presentation.

        Args:
            slide_layout (layout-from-template): The layout that is used.
        """
        slide = self.powerpoint_presentation.slides.add_slide(slide_layout)
        return slide

    def create_empty_slide(self):
        """
        create empty slide based on default empty template-slide
        """
        slide_layout = self.powerpoint_presentation.slide_layouts[0]
        self.add_slide(slide_layout)


# Main execution
if __name__ == "__main__":
    word_filename = "orde-van-dienst.docx"
    sermon = Sermon(word_filename)
    sermon.process_sermon()
